import streamlit as st
import pandas as pd
from pptx import Presentation
import io
from st_aggrid import AgGrid, GridOptionsBuilder

# 페이지 설정
st.set_page_config(page_title="문서 관리자", page_icon="📂", layout="wide")

st.title("📂 문서 통합 관리자")
st.markdown("""
엑셀, 파워포인트, PDF 문서를 업로드하여 내용을 확인하고 데이터를 추출할 수 있습니다.
""")

# 파일 업로더
uploaded_file = st.file_uploader(
    "문서를 업로드하세요", 
    type=['xlsx', 'csv', 'pptx', 'pdf'],
    help="지원 형식: Excel(.xlsx), CSV, PowerPoint(.pptx), PDF"
)

if uploaded_file:
    file_type = uploaded_file.name.split('.')[-1].lower()
    
    # ==========================================
    # 1. 엑셀/CSV 처리
    # ==========================================
    if file_type in ['xlsx', 'csv']:
        st.subheader(f"📊 엑셀/CSV 데이터: {uploaded_file.name}")
        
        try:
            if file_type == 'xlsx':
                df = pd.read_excel(uploaded_file)
            else:
                df = pd.read_csv(uploaded_file)
            
            # 탭으로 보기 방식 제공
            tab1, tab2 = st.tabs(["📋 데이터 그리드", "📈 데이터 요약"])
            
            with tab1:
                # AgGrid를 사용한 엑셀 같은 편집 UX
                gb = GridOptionsBuilder.from_dataframe(df)
                gb.configure_pagination(paginationAutoPageSize=True)
                gb.configure_side_bar()
                gb.configure_selection('multiple', use_checkbox=True, groupSelectsChildren="Group checkbox select children")
                gridOptions = gb.build()
                
                grid_response = AgGrid(
                    df,
                    gridOptions=gridOptions,
                    enable_enterprise_modules=False,
                    height=400,
                    fit_columns_on_grid_load=False,
                )
                
            with tab2:
                st.write("기초 통계량")
                st.dataframe(df.describe())
                
        except Exception as e:
            st.error(f"파일을 읽는 중 오류 발생: {e}")

    # ==========================================
    # 2. 파워포인트(PPTX) 처리
    # ==========================================
    elif file_type == 'pptx':
        st.subheader(f"📊 프레젠테이션 분석: {uploaded_file.name}")
        
        try:
            prs = Presentation(uploaded_file)
            
            st.info(f"총 슬라이드 수: {len(prs.slides)}장")
            
            # 슬라이드별 텍스트 추출
            slide_data = []
            for i, slide in enumerate(prs.slides):
                text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
                
                if text_content:
                    slide_data.append({"Slide": i+1, "Content": "\n".join(text_content)})
            
            if slide_data:
                df_slides = pd.DataFrame(slide_data)
                st.dataframe(df_slides, use_container_width=True)
                
                # 다운로드 버튼
                csv = df_slides.to_csv(index=False).encode('utf-8-sig')
                st.download_button(
                    label="텍스트 추출 결과 다운로드 (CSV)",
                    data=csv,
                    file_name=f"{uploaded_file.name}_extracted.csv",
                    mime="text/csv",
                )
            else:
                st.warning("텍스트를 추출할 수 없습니다 (이미지 위주 슬라이드 가능성).")
                
        except Exception as e:
            st.error(f"PPT 파일을 읽는 중 오류 발생: {e}")

    # ==========================================
    # 3. PDF 처리 (뷰어 및 텍스트)
    # ==========================================
    elif file_type == 'pdf':
        st.subheader(f"📄 PDF 문서: {uploaded_file.name}")
        
        # PDF 뷰어 (선택 사항)
        try:
            from streamlit_pdf_viewer import pdf_viewer
            binary_data = uploaded_file.getvalue()
            pdf_viewer(input=binary_data, width=700)
        except ImportError:
            st.warning("PDF 뷰어 라이브러리가 설치되지 않았습니다.")

else:
    st.info("👆 위 영역에 파일을 드래그하여 업로드하세요.")
    

    st.markdown("### 💡 활용 팁")
    st.markdown("""
    - **엑셀(Excel)**: 농사 일지, 비용 내역 등을 업로드하여 바로 표 형태로 확인하고 분석합니다.
    - **파워포인트(PPT)**: 기획서나 발표 자료를 올리면 텍스트만 추출하여 데이터베이스에 저장할 수 있습니다.
    - **PDF**: 매뉴얼이나 계약서를 업로드하여 내용을 바로 확인할 수 있습니다.
    """)

# ==========================================
# 4. 문서 생성 (Export) 데모
# ==========================================
st.divider()
st.subheader("📤 문서 생성 및 다운로드 (Export)")

col1, col2 = st.columns(2)

with col1:
    st.markdown("#### 📊 엑셀 리포트 생성")
    st.info("현재 화면의 데이터나 분석 결과를 엑셀 파일로 저장합니다.")
    
    if st.button("예시 엑셀 리포트 생성"):
        # 예시 데이터 생성
        data = {
            '작물': ['토마토', '오이', '딸기', '고추'],
            '면적(평)': [500, 300, 200, 400],
            '예상수익(만)': [1500, 900, 2000, 1200],
            '상태': ['재배중', '수확기', '파종전', '재배중']
        }
        df_export = pd.DataFrame(data)
        
        # 엑셀 생성
        output = io.BytesIO()
        with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
            df_export.to_excel(writer, sheet_name='작물현황', index=False)
            
            # 차트 추가 (선택사항)
            workbook = writer.book
            worksheet = writer.sheets['작물현황']
            chart = workbook.add_chart({'type': 'column'})
            chart.add_series({
                'values': '=작물현황!$C$2:$C$5',
                'categories': '=작물현황!$A$2:$A$5',
                'name': '예상수익',
            })
            worksheet.insert_chart('E2', chart)
            
        output.seek(0)
        
        st.download_button(
            label="📥 엑셀 리포트 다운로드",
            data=output,
            file_name="농장_현황_리포트.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

with col2:
    st.markdown("#### 📑 PPT 발표자료 생성")
    st.info("데이터를 기반으로 깔끔한 슬라이드를 자동 생성합니다.")
    
    ppt_title = st.text_input("프레젠테이션 제목", "2026년 영농 계획")
    
    if st.button("예시 PPT 생성"):
        prs = Presentation()
        
        # 제목 슬라이드
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = ppt_title
        subtitle.text = "Farm Calculator 자동 생성 리포트"
        
        # 내용 슬라이드
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "주요 목표"
        
        tf = body_shape.text_frame
        tf.text = "생산성 20% 향상"
        
        p = tf.add_paragraph()
        p.text = "데이터 기반 의사결정 도입"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "스마트팜 설비 자동화"
        p.level = 1
        
        # 저장
        ppt_output = io.BytesIO()
        prs.save(ppt_output)
        ppt_output.seek(0)
        
        st.download_button(
            label="📥 PPT 다운로드",
            data=ppt_output,
            file_name="영농계획.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

